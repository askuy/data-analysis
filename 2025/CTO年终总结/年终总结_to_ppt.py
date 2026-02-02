#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Convert 年终总结_ppt.html to PowerPoint format
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def add_title_slide(prs, title, subtitle, highlights):
    """Slide 1: Cover page with gradient background"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background (simulated with shape)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0x66, 0x7e, 0xea)
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.333), Inches(0.5))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
    p.alignment = PP_ALIGN.CENTER
    
    # Highlights
    x_start = 0.8
    card_width = 2.8
    gap = 0.3
    y = 4.2
    
    for i, (h_title, h_desc) in enumerate(highlights):
        x = x_start + i * (card_width + gap)
        
        # Card background
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(card_width), Inches(1.8))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0xff, 0xff, 0xff)
        card.fill.fore_color.brightness = 0.15
        card.line.fill.background()
        
        # Card title
        t_box = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.2), Inches(card_width - 0.3), Inches(0.4))
        tf = t_box.text_frame
        p = tf.paragraphs[0]
        p.text = h_title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
        
        # Card desc
        d_box = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.6), Inches(card_width - 0.3), Inches(1))
        tf = d_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = h_desc
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
    
    # Slide number
    add_slide_number(slide, "1", white=True)

def add_slide_number(slide, num, white=False):
    """Add slide number to bottom right"""
    num_box = slide.shapes.add_textbox(Inches(12.5), Inches(7), Inches(0.5), Inches(0.3))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(0xff, 0xff, 0xff) if white else RGBColor(0xad, 0xb5, 0xbd)
    p.alignment = PP_ALIGN.RIGHT

def add_header(slide, title, subtitle):
    """Add standard header to slide"""
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1a, 0x1a, 0x2e)
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(12.333), Inches(0.4))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0x6c, 0x75, 0x7d)

def add_card(slide, x, y, w, h, border_color=None):
    """Add a card shape"""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xff, 0xff, 0xff)
    card.line.color.rgb = RGBColor(0xe9, 0xec, 0xef)
    
    if border_color:
        border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.06), Inches(h))
        border.fill.solid()
        border.fill.fore_color.rgb = border_color
        border.line.fill.background()
    
    return card

def add_card_title(slide, x, y, text, color, icon=""):
    """Add card title with icon"""
    title_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(5.5), Inches(0.4))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{icon} {text}" if icon else text
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = color

def add_list_items(slide, x, y, items, marker="•", marker_color=None):
    """Add list items"""
    list_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(5.5), Inches(3))
    tf = list_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{marker} {item}"
        p.font.size = Pt(13)
        p.font.color.rgb = RGBColor(0x49, 0x50, 0x57)
        p.space_after = Pt(8)

def add_numbered_items(slide, x, y, items, badge_color=None):
    """Add numbered items with badges"""
    y_offset = y
    for i, (title, desc) in enumerate(items):
        # Number
        num_box = slide.shapes.add_textbox(Inches(x), Inches(y_offset), Inches(0.3), Inches(0.3))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(i + 1)
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = badge_color or RGBColor(0x66, 0x7e, 0xea)
        
        # Title
        t_box = slide.shapes.add_textbox(Inches(x + 0.35), Inches(y_offset - 0.05), Inches(5), Inches(0.3))
        tf = t_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x1a, 0x1a, 0x2e)
        
        # Description
        d_box = slide.shapes.add_textbox(Inches(x + 0.35), Inches(y_offset + 0.2), Inches(5), Inches(0.4))
        tf = d_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(0x6c, 0x75, 0x7d)
        
        y_offset += 0.55

def add_quote_box(slide, x, y, w, text):
    """Add quote box with gradient background"""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.6))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x66, 0x7e, 0xea)
    box.line.fill.background()
    
    text_box = slide.shapes.add_textbox(Inches(x), Inches(y + 0.15), Inches(w), Inches(0.4))
    tf = text_box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
    p.alignment = PP_ALIGN.CENTER

def add_result_cards(slide, y, items):
    """Add result cards at bottom"""
    card_width = 4.0
    x_start = 0.5
    gap = 0.166
    
    for i, (title, content) in enumerate(items):
        x = x_start + i * (card_width + gap)
        
        # Card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(card_width), Inches(1.2))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0xf8, 0xf9, 0xfa)
        card.line.color.rgb = RGBColor(0xe9, 0xec, 0xef)
        
        # Top border
        border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(card_width), Inches(0.04))
        border.fill.solid()
        border.fill.fore_color.rgb = RGBColor(0x34, 0x98, 0xdb)
        border.line.fill.background()
        
        # Title
        t_box = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.15), Inches(card_width - 0.3), Inches(0.3))
        tf = t_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x1a, 0x1a, 0x2e)
        
        # Content
        c_box = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.45), Inches(card_width - 0.3), Inches(0.7))
        tf = c_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = content
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(0x6c, 0x75, 0x7d)

def add_stat_cards(slide, x, y, stats):
    """Add statistics cards"""
    card_width = 3.8
    gap = 0.25
    
    for i, (value, label, sublabel) in enumerate(stats):
        cx = x + i * (card_width + gap)
        
        # Card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx), Inches(y), Inches(card_width), Inches(1.3))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0xf8, 0xf9, 0xfa)
        card.line.color.rgb = RGBColor(0xe9, 0xec, 0xef)
        
        # Value
        v_box = slide.shapes.add_textbox(Inches(cx), Inches(y + 0.2), Inches(card_width), Inches(0.6))
        tf = v_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x66, 0x7e, 0xea)
        p.alignment = PP_ALIGN.CENTER
        
        # Label
        l_box = slide.shapes.add_textbox(Inches(cx), Inches(y + 0.85), Inches(card_width), Inches(0.3))
        tf = l_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(0x6c, 0x75, 0x7d)
        p.alignment = PP_ALIGN.CENTER
        
        if sublabel:
            s_box = slide.shapes.add_textbox(Inches(cx), Inches(y + 1.05), Inches(card_width), Inches(0.2))
            tf = s_box.text_frame
            p = tf.paragraphs[0]
            p.text = sublabel
            p.font.size = Pt(10)
            p.font.color.rgb = RGBColor(0xad, 0xb5, 0xbd)
            p.alignment = PP_ALIGN.CENTER

def create_slide_2(prs):
    """Drive 的战略意义"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Drive 的战略意义", "愿景：文档届的 GitLab")
    
    # Past box
    past_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(5), Inches(2.8))
    past_box.fill.solid()
    past_box.fill.fore_color.rgb = RGBColor(0xfe, 0xf5, 0xf5)
    past_box.line.color.rgb = RGBColor(0xf8, 0xd7, 0xda)
    
    add_card_title(slide, 0.7, 1.65, "过去 (2024)", RGBColor(0xdc, 0x35, 0x45))
    add_list_items(slide, 0.7, 2.05, [
        "全站版本部署和运维复杂",
        "成本高、利润低",
        "定制化：小功能定制化",
        "广东电信发布需比照代码，耗时3天，SRE手动操作"
    ], "✗", RGBColor(0xe7, 0x4c, 0x3c))
    
    # Arrow
    arrow_box = slide.shapes.add_textbox(Inches(5.7), Inches(2.5), Inches(0.8), Inches(0.5))
    tf = arrow_box.text_frame
    p = tf.paragraphs[0]
    p.text = "→"
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(0x66, 0x7e, 0xea)
    
    # Now box
    now_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(1.5), Inches(6), Inches(2.8))
    now_box.fill.solid()
    now_box.fill.fore_color.rgb = RGBColor(0xf0, 0xfd, 0xf4)
    now_box.line.color.rgb = RGBColor(0xd1, 0xfa, 0xe5)
    
    add_card_title(slide, 6.7, 1.65, "现在 (2025)", RGBColor(0x10, 0xb9, 0x81))
    add_list_items(slide, 6.7, 2.05, [
        "Drive 简化部署",
        "成本低、利润可观",
        "定制化：大功能定制化",
        "直接拿包，SRE、研发均可操作"
    ], "✓", RGBColor(0x27, 0xae, 0x60))
    
    # Result cards
    add_result_cards(slide, 4.5, [
        ("🏢 广东电信", "数据迁移：6-8小时 → 30分钟"),
        ("📚 好未来", "灰度迁移方案落地"),
        ("📱 OPPO", "双环境迁移方案")
    ])
    
    add_slide_number(slide, "2")

def create_slide_3(prs):
    """私有化交付"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "私有化交付", "在 TOP10 客户上投入大量工时，确保交付质量")
    
    # Stats
    add_stat_cards(slide, 0.5, 1.5, [
        ("17", "大版本升级", ""),
        ("13", "正式交付", ""),
        ("15", "POC", "")
    ])
    
    # Left card
    add_card(slide, 0.5, 3.1, 6, 3.3, RGBColor(0x34, 0x98, 0xdb))
    add_card_title(slide, 0.7, 3.25, "自动化发布和升级", RGBColor(0x34, 0x98, 0xdb), "🚀")
    add_list_items(slide, 0.7, 3.7, [
        "确保 Drive 交付升级顺畅",
        "日志系统：客户侧问题尽快发现",
        "工具可视化：极大降低运维成本",
        "POC 部署变得轻量化、去专业化"
    ], "✓", RGBColor(0x27, 0xae, 0x60))
    
    # Right card
    add_card(slide, 6.833, 3.1, 6, 3.3, RGBColor(0x9b, 0x59, 0xb6))
    add_card_title(slide, 7.05, 3.25, "技术突破", RGBColor(0x9b, 0x59, 0xb6), "⚡")
    add_numbered_items(slide, 7.1, 3.75, [
        ("应用表格", "百万单元格：无法打开 → 4.5s | 高级权限系统上线"),
        ("极速 SDK", "前端入口统一，性能大幅提升")
    ], RGBColor(0x9b, 0x59, 0xb6))
    
    add_slide_number(slide, "3")

def create_slide_4(prs):
    """稳定性"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "稳定性", "Bug 数逐月下降趋势")
    
    # Left - Table card
    add_card(slide, 0.5, 1.5, 6, 4.5)
    add_card_title(slide, 0.7, 1.65, "📊 月度缺陷数据", RGBColor(0x1a, 0x1a, 0x2e))
    
    # Table header
    table_y = 2.1
    headers = ["月份", "SaaS", "私有化", "趋势"]
    header_box = slide.shapes.add_textbox(Inches(0.7), Inches(table_y), Inches(5.5), Inches(0.35))
    tf = header_box.text_frame
    p = tf.paragraphs[0]
    p.text = "月份      SaaS      私有化      趋势"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x49, 0x50, 0x57)
    
    # Table data
    data = [
        ("1月", "162", "33", "▼", True),
        ("2月", "73", "95", "", False),
        ("3月", "101", "84", "", False),
        ("6月", "65", "73", "", False),
        ("9月", "41", "22", "▼", False),
        ("12月", "59", "50", "▼", False),
    ]
    
    row_y = table_y + 0.4
    for month, saas, priv, trend, is_red in data:
        row_box = slide.shapes.add_textbox(Inches(0.7), Inches(row_y), Inches(5.5), Inches(0.3))
        tf = row_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{month}        {saas}         {priv}            {trend}"
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(0x49, 0x50, 0x57)
        row_y += 0.35
    
    # Right - Stats cards
    add_card(slide, 6.833, 1.5, 6, 2.2, RGBColor(0x27, 0xae, 0x60))
    add_card_title(slide, 7.05, 1.65, "关键指标", RGBColor(0x27, 0xae, 0x60), "✓")
    
    # Mini stat cards
    stat_box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.1), Inches(2.1), Inches(2.6), Inches(1.3))
    stat_box1.fill.solid()
    stat_box1.fill.fore_color.rgb = RGBColor(0xf8, 0xf9, 0xfa)
    stat_box1.line.color.rgb = RGBColor(0xe9, 0xec, 0xef)
    
    v_box = slide.shapes.add_textbox(Inches(7.1), Inches(2.25), Inches(2.6), Inches(0.5))
    tf = v_box.text_frame
    p = tf.paragraphs[0]
    p.text = "99.99%"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x66, 0x7e, 0xea)
    p.alignment = PP_ALIGN.CENTER
    
    l_box = slide.shapes.add_textbox(Inches(7.1), Inches(2.8), Inches(2.6), Inches(0.3))
    tf = l_box.text_frame
    p = tf.paragraphs[0]
    p.text = "SaaS 接口稳定性"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(0x6c, 0x75, 0x7d)
    p.alignment = PP_ALIGN.CENTER
    
    stat_box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10), Inches(2.1), Inches(2.6), Inches(1.3))
    stat_box2.fill.solid()
    stat_box2.fill.fore_color.rgb = RGBColor(0xf8, 0xf9, 0xfa)
    stat_box2.line.color.rgb = RGBColor(0xe9, 0xec, 0xef)
    
    v_box = slide.shapes.add_textbox(Inches(10), Inches(2.25), Inches(2.6), Inches(0.5))
    tf = v_box.text_frame
    p = tf.paragraphs[0]
    p.text = "99.9%"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x66, 0x7e, 0xea)
    p.alignment = PP_ALIGN.CENTER
    
    l_box = slide.shapes.add_textbox(Inches(10), Inches(2.8), Inches(2.6), Inches(0.5))
    tf = l_box.text_frame
    p = tf.paragraphs[0]
    p.text = "导入导出稳定性\n99.1% → 99.9%"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(0x6c, 0x75, 0x7d)
    p.alignment = PP_ALIGN.CENTER
    
    # Bug types card
    add_card(slide, 6.833, 3.9, 6, 2.1, RGBColor(0xe6, 0x7e, 0x22))
    add_card_title(slide, 7.05, 4.05, "SaaS 缺陷类型 Top3", RGBColor(0xe6, 0x7e, 0x22), "⚠")
    add_list_items(slide, 7.1, 4.5, [
        "代码问题：384（69.3%）",
        "非问题：62",
        "需求上线引入：41"
    ])
    
    add_slide_number(slide, "4")

def create_slide_5(prs):
    """稳定性保障措施"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "稳定性保障措施", "从被客户问题拖累，到主动防御")
    
    # Left card
    add_card(slide, 0.5, 1.5, 6, 4.2, RGBColor(0x34, 0x98, 0xdb))
    add_card_title(slide, 0.7, 1.65, "SaaS 保障", RGBColor(0x34, 0x98, 0xdb), "🛡️")
    add_numbered_items(slide, 0.75, 2.15, [
        ("值班制度", "7x24 小时响应机制"),
        ("报警配置", "150 条报警规则 | 重构后业务 > 0 即报警"),
        ("导入导出优化", "大量 SOP 支撑 | 工单减少 | 人员可更多做需求")
    ], RGBColor(0x34, 0x98, 0xdb))
    
    # Right card
    add_card(slide, 6.833, 1.5, 6, 4.2, RGBColor(0x9b, 0x59, 0xb6))
    add_card_title(slide, 7.05, 1.65, "长期质量策略", RGBColor(0x9b, 0x59, 0xb6), "🎯")
    
    # Short term box
    short_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.1), Inches(2.15), Inches(5.5), Inches(0.9))
    short_box.fill.solid()
    short_box.fill.fore_color.rgb = RGBColor(0xfe, 0xf5, 0xf5)
    short_box.line.color.rgb = RGBColor(0xf8, 0xd7, 0xda)
    
    t_box = slide.shapes.add_textbox(Inches(7.25), Inches(2.25), Inches(5.2), Inches(0.3))
    tf = t_box.text_frame
    p = tf.paragraphs[0]
    p.text = "短期：交付团队 + SOP"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xdc, 0x35, 0x45)
    
    d_box = slide.shapes.add_textbox(Inches(7.25), Inches(2.55), Inches(5.2), Inches(0.3))
    tf = d_box.text_frame
    p = tf.paragraphs[0]
    p.text = "快速响应客户问题"
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(0x6c, 0x75, 0x7d)
    
    # Long term box
    long_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.1), Inches(3.2), Inches(5.5), Inches(2.3))
    long_box.fill.solid()
    long_box.fill.fore_color.rgb = RGBColor(0xf0, 0xfd, 0xf4)
    long_box.line.color.rgb = RGBColor(0xd1, 0xfa, 0xe5)
    
    t_box = slide.shapes.add_textbox(Inches(7.25), Inches(3.3), Inches(5.2), Inches(0.3))
    tf = t_box.text_frame
    p = tf.paragraphs[0]
    p.text = "长期：质量稳定"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x10, 0xb9, 0x81)
    
    add_list_items(slide, 7.25, 3.65, [
        "单元测试（AI 辅助）",
        "自动化测试（AI 辅助）",
        "运行错误：SaaS Error 日志 > 0 即报警"
    ], "✓", RGBColor(0x27, 0xae, 0x60))
    
    # Quote
    add_quote_box(slide, 0.5, 5.9, 12.333, "目标：客户不想报 Bug，而是我们主动发现并解决问题")
    
    add_slide_number(slide, "5")

def create_slide_6(prs):
    """内部体验机制（Dogfooding）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "内部体验机制（Dogfooding）", '从"被动修补"到"主动优化"的思维转变')
    
    # Left card - Pain points
    add_card(slide, 0.5, 1.5, 6, 2.5, RGBColor(0xe7, 0x4c, 0x3c))
    add_card_title(slide, 0.7, 1.65, "核心痛点：脱离客户视角", RGBColor(0xe7, 0x4c, 0x3c), "⚠")
    add_list_items(slide, 0.7, 2.1, [
        "shimo.im 与 SDK 产品架构不统一，未真正体验私有化产品",
        "研发团队缺乏对私有化全站和 SDK 产品的真实使用体验",
        "对客户反馈的痛点重视不足，响应滞后"
    ], "✗", RGBColor(0xe7, 0x4c, 0x3c))
    
    # Right card - Solutions
    add_card(slide, 6.833, 1.5, 6, 2.5, RGBColor(0x27, 0xae, 0x60))
    add_card_title(slide, 7.05, 1.65, "变革举措", RGBColor(0x27, 0xae, 0x60), "✓")
    add_numbered_items(slide, 7.1, 2.1, [
        ("引入 Drive 团队", '作为 SDK 首个"内部客户"，在开发阶段即暴露集成与体验问题'),
        ("构建 shimo.net", "全员使用 Drive + SDK 产品，将私有化升级作为内部日常办公环境"),
        ("真实环境暴露", "通过私有化升级方式处理 shimo.net，把问题暴露到内部使用环境")
    ], RGBColor(0x27, 0xae, 0x60))
    
    # Result cards
    add_result_cards(slide, 4.2, [
        ("💡 关键发现", "客户提的问题都是对的\n例：协作头像 API 复杂、搜索对接 API 复杂"),
        ("🎯 架构对齐", "与客户环境一致才能解决问题\n更早发现问题，在交付前拦截缺陷"),
        ("🚀 主动优化", "有些问题不应属于定制化\n研发主动处理，提前优化体验")
    ])
    
    # Quote
    add_quote_box(slide, 0.5, 5.6, 12.333, "客户提的问题都是对的 —— 只有自己用了才知道")
    
    add_slide_number(slide, "6")

def create_slide_7(prs):
    """AI 提效"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "AI 提效", "使用 AI 不是简单与 AI 对话 —— Skill、MCP、Agent")
    
    # Left card
    add_card(slide, 0.5, 1.5, 6, 5, RGBColor(0x9b, 0x59, 0xb6))
    add_card_title(slide, 0.7, 1.65, "AI 应用场景", RGBColor(0x9b, 0x59, 0xb6), "🤖")
    add_numbered_items(slide, 0.75, 2.15, [
        ("海外类飞书项目", "1 个月完成 | 古法编程不可取，先用 AI 试试"),
        ("前后端联调", "AI 自己登录、探测浏览器报错、检测数据库数据准确性"),
        ("项目预估", "给 CSV，AI 直接预估大项目、定制化项目工时"),
        ("测试用例", "直接用 AI 转成 Playwright 自动化测试")
    ], RGBColor(0x9b, 0x59, 0xb6))
    
    # Right card
    add_card(slide, 6.833, 1.5, 6, 3.2, RGBColor(0x34, 0x98, 0xdb))
    add_card_title(slide, 7.05, 1.65, "Word 预览突破", RGBColor(0x34, 0x98, 0xdb), "📄")
    
    # Feature box
    feature_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.1), Inches(2.15), Inches(5.5), Inches(2.3))
    feature_box.fill.solid()
    feature_box.fill.fore_color.rgb = RGBColor(0xf0, 0xfd, 0xf4)
    feature_box.line.color.rgb = RGBColor(0xd1, 0xfa, 0xe5)
    
    t_box = slide.shapes.add_textbox(Inches(7.25), Inches(2.25), Inches(5.2), Inches(0.3))
    tf = t_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Opus + MCP"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x10, 0xb9, 0x81)
    
    add_list_items(slide, 7.25, 2.6, [
        "直接理解 Word OOXML",
        "快速提升新版 Word 预览",
        "比老石墨增加更多排版能力",
        "准确性大幅提升"
    ], "✓", RGBColor(0x27, 0xae, 0x60))
    
    # Quote
    quote_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.833), Inches(4.9), Inches(6), Inches(1))
    quote_box.fill.solid()
    quote_box.fill.fore_color.rgb = RGBColor(0x1a, 0x1a, 0x2e)
    quote_box.line.fill.background()
    
    q_text = slide.shapes.add_textbox(Inches(6.833), Inches(5.15), Inches(6), Inches(0.6))
    tf = q_text.text_frame
    p = tf.paragraphs[0]
    p.text = "古法编程不可取\n无论如何先要用 AI 试试"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
    p.alignment = PP_ALIGN.CENTER
    
    add_slide_number(slide, "7")

def create_slide_8(prs):
    """研发架构"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "研发架构", "Drive + SDK + SDK Core 三层架构")
    
    # Architecture diagram
    arch_y = 1.8
    box_width = 3.2
    gap = 0.8
    
    # Drive box
    drive_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(arch_y), Inches(box_width), Inches(1.5))
    drive_box.fill.solid()
    drive_box.fill.fore_color.rgb = RGBColor(0x66, 0x7e, 0xea)
    drive_box.line.fill.background()
    
    d_text = slide.shapes.add_textbox(Inches(2), Inches(arch_y + 0.35), Inches(box_width), Inches(0.5))
    tf = d_text.text_frame
    p = tf.paragraphs[0]
    p.text = "Drive"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
    p.alignment = PP_ALIGN.CENTER
    
    d_sub = slide.shapes.add_textbox(Inches(2), Inches(arch_y + 0.85), Inches(box_width), Inches(0.3))
    tf = d_sub.text_frame
    p = tf.paragraphs[0]
    p.text = "最外层"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
    p.alignment = PP_ALIGN.CENTER
    
    # Arrow 1
    arrow1 = slide.shapes.add_textbox(Inches(5.3), Inches(arch_y + 0.45), Inches(0.5), Inches(0.5))
    tf = arrow1.text_frame
    p = tf.paragraphs[0]
    p.text = "→"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0x66, 0x7e, 0xea)
    
    # SDK box
    sdk_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.9), Inches(arch_y), Inches(box_width), Inches(1.5))
    sdk_box.fill.solid()
    sdk_box.fill.fore_color.rgb = RGBColor(0x34, 0x98, 0xdb)
    sdk_box.line.fill.background()
    
    s_text = slide.shapes.add_textbox(Inches(5.9), Inches(arch_y + 0.35), Inches(box_width), Inches(0.5))
    tf = s_text.text_frame
    p = tf.paragraphs[0]
    p.text = "SDK"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
    p.alignment = PP_ALIGN.CENTER
    
    s_sub = slide.shapes.add_textbox(Inches(5.9), Inches(arch_y + 0.85), Inches(box_width), Inches(0.3))
    tf = s_sub.text_frame
    p = tf.paragraphs[0]
    p.text = "胶水层"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
    p.alignment = PP_ALIGN.CENTER
    
    # Arrow 2
    arrow2 = slide.shapes.add_textbox(Inches(9.2), Inches(arch_y + 0.45), Inches(0.5), Inches(0.5))
    tf = arrow2.text_frame
    p = tf.paragraphs[0]
    p.text = "→"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0x66, 0x7e, 0xea)
    
    # SDK Core box
    core_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.8), Inches(arch_y), Inches(box_width), Inches(1.5))
    core_box.fill.solid()
    core_box.fill.fore_color.rgb = RGBColor(0x27, 0xae, 0x60)
    core_box.line.fill.background()
    
    c_text = slide.shapes.add_textbox(Inches(9.8), Inches(arch_y + 0.35), Inches(box_width), Inches(0.5))
    tf = c_text.text_frame
    p = tf.paragraphs[0]
    p.text = "SDK Core"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
    p.alignment = PP_ALIGN.CENTER
    
    c_sub = slide.shapes.add_textbox(Inches(9.8), Inches(arch_y + 0.85), Inches(box_width), Inches(0.3))
    tf = c_sub.text_frame
    p = tf.paragraphs[0]
    p.text = "纯套件"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
    p.alignment = PP_ALIGN.CENTER
    
    # Description table
    table_data = [
        ("Drive", "职责：最外层，负责产品体验和商业化", "优势：聚焦核心场景和卖点"),
        ("SDK", "职责：胶水层，统一前后端逻辑", "优势：统一私有化对接团队"),
        ("SDK Core", "职责：纯套件，专注 UI 库和 JS API", "优势：打造高可用 JS SDK API"),
    ]
    
    row_y = 3.6
    for name, resp, adv in table_data:
        row_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(row_y), Inches(12.333), Inches(0.9))
        row_box.fill.solid()
        row_box.fill.fore_color.rgb = RGBColor(0xf8, 0xf9, 0xfa)
        row_box.line.fill.background()
        
        n_box = slide.shapes.add_textbox(Inches(0.7), Inches(row_y + 0.3), Inches(1.5), Inches(0.4))
        tf = n_box.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x66, 0x7e, 0xea)
        
        r_box = slide.shapes.add_textbox(Inches(2.3), Inches(row_y + 0.3), Inches(5), Inches(0.4))
        tf = r_box.text_frame
        p = tf.paragraphs[0]
        p.text = resp
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(0x49, 0x50, 0x57)
        
        a_box = slide.shapes.add_textbox(Inches(7.5), Inches(row_y + 0.3), Inches(5), Inches(0.4))
        tf = a_box.text_frame
        p = tf.paragraphs[0]
        p.text = adv
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(0x49, 0x50, 0x57)
        
        row_y += 1.05
    
    add_slide_number(slide, "8")

def create_slide_9(prs):
    """问题与挑战"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "问题与挑战", "需要反思的问题 & 持续面临的挑战")
    
    # Left card - Problems
    add_card(slide, 0.5, 1.5, 6, 4.5, RGBColor(0xe7, 0x4c, 0x3c))
    add_card_title(slide, 0.7, 1.65, "需要反思的问题", RGBColor(0xe7, 0x4c, 0x3c), "🔍")
    add_numbered_items(slide, 0.75, 2.15, [
        ("客户反馈重视不够", "Dogfooding 后验证：客户是对的"),
        ("代码问题占比高", "Bug 的 69.3% 为代码问题"),
        ("人员单点风险", "每个研发仅负责一个套件")
    ], RGBColor(0xe7, 0x4c, 0x3c))
    
    # Right card - Challenges
    add_card(slide, 6.833, 1.5, 6, 4.5, RGBColor(0xe6, 0x7e, 0x22))
    add_card_title(slide, 7.05, 1.65, "持续面临的挑战", RGBColor(0xe6, 0x7e, 0x22), "⚠")
    add_numbered_items(slide, 7.1, 2.15, [
        ("专业文档功能不齐", "功能覆盖度仍需提升"),
        ("幻灯片 H5 编辑缺失", "移动端不支持编辑"),
        ("交付持续消耗", "占总工时 20%+")
    ], RGBColor(0xe6, 0x7e, 0x22))
    
    add_slide_number(slide, "9")

def create_slide_10(prs):
    """Final slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0x1a, 0x1a, 0x2e)
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "以客户为中心，迈向新征程"
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
    p.alignment = PP_ALIGN.CENTER
    
    # Vision
    vision_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.4), Inches(12.333), Inches(0.5))
    tf = vision_box.text_frame
    p = tf.paragraphs[0]
    p.text = '2026：向着"文档界的 GitLab"坚定迈进'
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0x66, 0x7e, 0xea)
    p.alignment = PP_ALIGN.CENTER
    
    # Summary items
    summary_items = [
        ("🚀 Drive 战略落地", "验证架构演进"),
        ("👥 体验机制建立", "回归客户视角"),
        ("⚙️ 自动化部署", "赋能商业推广"),
        ("📦 产品增强", "补齐套件功能短板\nSDK/JS API 全面增强")
    ]
    
    x_start = 0.8
    card_width = 2.9
    gap = 0.25
    y = 3.2
    
    for i, (title, desc) in enumerate(summary_items):
        x = x_start + i * (card_width + gap)
        
        # Card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(card_width), Inches(1.8))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0xff, 0xff, 0xff)
        card.fill.fore_color.brightness = 0.1
        card.line.fill.background()
        
        # Title
        t_box = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.25), Inches(card_width - 0.3), Inches(0.4))
        tf = t_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x66, 0x7e, 0xea)
        
        # Desc
        d_box = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.7), Inches(card_width - 0.3), Inches(0.9))
        tf = d_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
    
    # Bottom stats
    stats = [
        ("99.99%", "SaaS 接口稳定性", RGBColor(0x66, 0x7e, 0xea)),
        ("99.9%", "导入导出稳定性", RGBColor(0x27, 0xae, 0x60)),
        ("Bug ↓", "逐月下降趋势", RGBColor(0xe6, 0x7e, 0x22))
    ]
    
    stat_x_start = 3
    stat_gap = 2.8
    stat_y = 5.5
    
    for i, (value, label, color) in enumerate(stats):
        sx = stat_x_start + i * stat_gap
        
        v_box = slide.shapes.add_textbox(Inches(sx), Inches(stat_y), Inches(2.5), Inches(0.6))
        tf = v_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER
        
        l_box = slide.shapes.add_textbox(Inches(sx), Inches(stat_y + 0.55), Inches(2.5), Inches(0.3))
        tf = l_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
        p.alignment = PP_ALIGN.CENTER
    
    add_slide_number(slide, "10", white=True)

def main():
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Cover
    add_title_slide(prs, "2025年度研发中心工作汇报", "以客户为中心，迈向新征程", [
        ("Drive 战略落地", "从零组建团队\n半年交付三大标杆客户"),
        ("内部体验机制", "构建 shimo.net\n回归客户视角"),
        ("自动化部署", "POC 无需立项\n非 SRE 亦可操作"),
        ("AI 提效", "一个月完成类飞书产品\n全面赋能研发流程")
    ])
    
    # Slide 2: Drive Strategic Significance
    create_slide_2(prs)
    
    # Slide 3: Delivery
    create_slide_3(prs)
    
    # Slide 4: Stability
    create_slide_4(prs)
    
    # Slide 5: Stability Measures
    create_slide_5(prs)
    
    # Slide 6: Dogfooding
    create_slide_6(prs)
    
    # Slide 7: AI Enhancement
    create_slide_7(prs)
    
    # Slide 8: Architecture
    create_slide_8(prs)
    
    # Slide 9: Challenges
    create_slide_9(prs)
    
    # Slide 10: Final
    create_slide_10(prs)
    
    # Save
    output_path = "年终总结.pptx"
    prs.save(output_path)
    print(f"PPT saved to: {output_path}")

if __name__ == "__main__":
    main()
