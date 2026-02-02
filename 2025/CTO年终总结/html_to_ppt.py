#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Convert dogfooding HTML slide to PowerPoint format
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml

def create_dogfooding_ppt():
    # Create presentation with 16:9 aspect ratio
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Add blank slide
    blank_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_layout)
    
    # Define colors
    title_color = RGBColor(0x1a, 0x1a, 0x2e)
    subtitle_color = RGBColor(0x6c, 0x75, 0x7d)
    pain_color = RGBColor(0xe7, 0x4c, 0x3c)
    solution_color = RGBColor(0x27, 0xae, 0x60)
    text_color = RGBColor(0x49, 0x50, 0x57)
    
    # ========== Title Section ==========
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.6))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "内部体验机制（Dogfooding）"
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = title_color
    title_para.alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(12.333), Inches(0.4))
    subtitle_frame = subtitle_box.text_frame
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.text = '从"被动修补"到"主动优化"的思维转变'
    subtitle_para.font.size = Pt(20)
    subtitle_para.font.color.rgb = subtitle_color
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # ========== Left Card - Pain Points ==========
    # Card background
    pain_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.5),
        Inches(6), Inches(2.8)
    )
    pain_card.fill.solid()
    pain_card.fill.fore_color.rgb = RGBColor(0xff, 0xff, 0xff)
    pain_card.line.color.rgb = RGBColor(0xe9, 0xec, 0xef)
    
    # Left border effect (simulated with thin rectangle)
    pain_border = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(1.5),
        Inches(0.06), Inches(2.8)
    )
    pain_border.fill.solid()
    pain_border.fill.fore_color.rgb = pain_color
    pain_border.line.fill.background()
    
    # Pain title
    pain_title = slide.shapes.add_textbox(Inches(0.7), Inches(1.65), Inches(5.6), Inches(0.4))
    pain_title_frame = pain_title.text_frame
    pain_title_para = pain_title_frame.paragraphs[0]
    pain_title_para.text = "⚠ 核心痛点：脱离客户视角"
    pain_title_para.font.size = Pt(18)
    pain_title_para.font.bold = True
    pain_title_para.font.color.rgb = pain_color
    
    # Pain list items
    pain_items = [
        "shimo.im 与 SDK 产品架构不统一，未真正体验私有化产品",
        "研发团队缺乏对私有化全站和 SDK 产品的真实使用体验",
        "对客户反馈的痛点重视不足，响应滞后"
    ]
    
    pain_list = slide.shapes.add_textbox(Inches(0.8), Inches(2.15), Inches(5.5), Inches(2.0))
    pain_list_frame = pain_list.text_frame
    pain_list_frame.word_wrap = True
    
    for i, item in enumerate(pain_items):
        if i == 0:
            para = pain_list_frame.paragraphs[0]
        else:
            para = pain_list_frame.add_paragraph()
        para.text = "✗ " + item
        para.font.size = Pt(14)
        para.font.color.rgb = text_color
        para.space_after = Pt(12)
    
    # ========== Right Card - Solutions ==========
    solution_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.833), Inches(1.5),
        Inches(6), Inches(2.8)
    )
    solution_card.fill.solid()
    solution_card.fill.fore_color.rgb = RGBColor(0xff, 0xff, 0xff)
    solution_card.line.color.rgb = RGBColor(0xe9, 0xec, 0xef)
    
    # Left border effect
    solution_border = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(6.833), Inches(1.5),
        Inches(0.06), Inches(2.8)
    )
    solution_border.fill.solid()
    solution_border.fill.fore_color.rgb = solution_color
    solution_border.line.fill.background()
    
    # Solution title
    solution_title = slide.shapes.add_textbox(Inches(7.05), Inches(1.65), Inches(5.6), Inches(0.4))
    solution_title_frame = solution_title.text_frame
    solution_title_para = solution_title_frame.paragraphs[0]
    solution_title_para.text = "✓ 变革举措"
    solution_title_para.font.size = Pt(18)
    solution_title_para.font.bold = True
    solution_title_para.font.color.rgb = solution_color
    
    # Solution items
    solutions = [
        ("引入 Drive 团队", '作为 SDK 首个"内部客户"，在开发阶段即暴露集成与体验问题'),
        ("构建 shimo.net", "全员使用 Drive + SDK 产品，将私有化升级作为内部日常办公环境"),
        ("真实环境暴露", "通过私有化升级方式处理 shimo.net，把问题暴露到内部使用环境")
    ]
    
    y_offset = 2.15
    for i, (title, desc) in enumerate(solutions):
        # Number circle (simulated with text)
        num_box = slide.shapes.add_textbox(Inches(7.1), Inches(y_offset), Inches(0.3), Inches(0.3))
        num_frame = num_box.text_frame
        num_para = num_frame.paragraphs[0]
        num_para.text = str(i + 1)
        num_para.font.size = Pt(12)
        num_para.font.bold = True
        num_para.font.color.rgb = solution_color
        
        # Solution title
        sol_title_box = slide.shapes.add_textbox(Inches(7.45), Inches(y_offset - 0.05), Inches(5.2), Inches(0.3))
        sol_title_frame = sol_title_box.text_frame
        sol_title_para = sol_title_frame.paragraphs[0]
        sol_title_para.text = title
        sol_title_para.font.size = Pt(14)
        sol_title_para.font.bold = True
        sol_title_para.font.color.rgb = title_color
        
        # Solution description
        sol_desc_box = slide.shapes.add_textbox(Inches(7.45), Inches(y_offset + 0.25), Inches(5.2), Inches(0.5))
        sol_desc_frame = sol_desc_box.text_frame
        sol_desc_frame.word_wrap = True
        sol_desc_para = sol_desc_frame.paragraphs[0]
        sol_desc_para.text = desc
        sol_desc_para.font.size = Pt(12)
        sol_desc_para.font.color.rgb = subtitle_color
        
        y_offset += 0.7
    
    # ========== Bottom Section - Results ==========
    results = [
        ("💡 关键发现", "客户提的问题都是对的\n例：协作头像 API 复杂、搜索对接 API 复杂"),
        ("🎯 架构对齐", "与客户环境一致才能解决问题\n更早发现问题，在交付前拦截缺陷"),
        ("🚀 主动优化", "有些问题不应属于定制化\n研发主动处理，提前优化体验")
    ]
    
    card_width = 4.0
    card_start_x = 0.5
    card_y = 4.5
    
    for i, (title, content) in enumerate(results):
        x_pos = card_start_x + i * (card_width + 0.166)
        
        # Result card background
        result_card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x_pos), Inches(card_y),
            Inches(card_width), Inches(1.4)
        )
        result_card.fill.solid()
        result_card.fill.fore_color.rgb = RGBColor(0xf8, 0xf9, 0xfa)
        result_card.line.color.rgb = RGBColor(0xe9, 0xec, 0xef)
        
        # Top border (gradient effect simulated)
        top_border = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x_pos), Inches(card_y),
            Inches(card_width), Inches(0.05)
        )
        top_border.fill.solid()
        top_border.fill.fore_color.rgb = RGBColor(0x34, 0x98, 0xdb)
        top_border.line.fill.background()
        
        # Result title
        result_title_box = slide.shapes.add_textbox(
            Inches(x_pos + 0.15), Inches(card_y + 0.15),
            Inches(card_width - 0.3), Inches(0.35)
        )
        result_title_frame = result_title_box.text_frame
        result_title_para = result_title_frame.paragraphs[0]
        result_title_para.text = title
        result_title_para.font.size = Pt(14)
        result_title_para.font.bold = True
        result_title_para.font.color.rgb = title_color
        
        # Result content
        result_content_box = slide.shapes.add_textbox(
            Inches(x_pos + 0.15), Inches(card_y + 0.5),
            Inches(card_width - 0.3), Inches(0.8)
        )
        result_content_frame = result_content_box.text_frame
        result_content_frame.word_wrap = True
        result_content_para = result_content_frame.paragraphs[0]
        result_content_para.text = content
        result_content_para.font.size = Pt(11)
        result_content_para.font.color.rgb = subtitle_color
    
    # ========== Quote Box ==========
    quote_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(6.1),
        Inches(12.333), Inches(0.7)
    )
    quote_box.fill.solid()
    quote_box.fill.fore_color.rgb = RGBColor(0x66, 0x7e, 0xea)
    quote_box.line.fill.background()
    
    quote_text = slide.shapes.add_textbox(Inches(0.5), Inches(6.25), Inches(12.333), Inches(0.4))
    quote_frame = quote_text.text_frame
    quote_para = quote_frame.paragraphs[0]
    quote_para.text = '"客户提的问题都是对的 —— 只有自己用了才知道"'
    quote_para.font.size = Pt(18)
    quote_para.font.bold = True
    quote_para.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
    quote_para.alignment = PP_ALIGN.CENTER
    
    # Save the presentation
    output_path = "dogfooding.pptx"
    prs.save(output_path)
    print(f"PPT saved to: {output_path}")
    return output_path

if __name__ == "__main__":
    create_dogfooding_ppt()
